import typing
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER


def pptx_to_text_pages(f: typing.BinaryIO, use_form_reco: bool = False) -> list[str]:
    """
    Extracts and converts text, tables, charts, and grouped shapes from a PPTX file into Markdown format.
    """
    prs = Presentation(f)
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_content = [f"Slide {slide_num}"]
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    slide_content.extend(handle_text_elements(shape))

                if shape.has_table:
                    slide_content.extend(handle_tables(shape))

                if shape.has_chart:
                    slide_content.extend(handle_charts(shape))

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    slide_content.extend(handle_grouped_shapes(shape))

                # if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                #     slide_content.extend(handle_pictures(shape))

            except Exception as e:
                slide_content.append(f"  Error processing shape: {e}")

        if slide.has_notes_slide:
            slide_content.extend(handle_author_notes(slide))

        slides_text.append("\n".join(slide_content) + "\n")

    return slides_text


def handle_text_elements(shape) -> list[str]:
    """
    Handles text elements within a shape, including lists.
    """
    text_elements = []
    namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    current_list_type = None
    list_item_index = 0

    for paragraph in shape.text_frame.paragraphs:
        p = paragraph._element
        paragraph_text = ""
        is_list_item = False

        # Determine list type
        if p.find(".//a:buChar", namespaces=namespaces) is not None:
            current_list_type = "Bullet"
            is_list_item = True
        elif p.find(".//a:buAutoNum", namespaces=namespaces) is not None:
            current_list_type = "Numbered"
            is_list_item = True
        elif paragraph.level > 0:  # Indented text is also treated as a list
            current_list_type = "Bullet"
            is_list_item = True
        else:
            current_list_type = None
            list_item_index = 0  # Reset numbering if no list

        # Process paragraph text
        for run in p.iterfind(".//a:r", namespaces=namespaces):
            run_text = run.text.strip() if run.text else ""
            if run_text:
                paragraph_text += run_text

        if is_list_item:
            if current_list_type == "Numbered":
                list_item_index += 1
                list_prefix = f"{list_item_index}."
            else:
                list_prefix = "•"  # Default bullet symbol
            text_elements.append(f"{list_prefix} {paragraph_text}")
        else:
            # Handle placeholders for titles or subtitles
            if shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                if placeholder_type == PP_PLACEHOLDER.TITLE:
                    text_elements.append(f"TITLE: {paragraph_text}")
                elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                    text_elements.append(f"SECTION_HEADER: {paragraph_text}")
                else:
                    text_elements.append(paragraph_text)
            else:
                text_elements.append(paragraph_text)

    return text_elements


def handle_tables(shape) -> list[str]:
    """
    Handles tables within a shape, converting them into Markdown format.
    """

    if not hasattr(shape, "has_table") or not shape.has_table:
        return []
    table = shape.table
    table_xml = shape._element

    num_rows = len(table.rows)
    num_cols = len(table.columns)
    grid = [["" for _ in range(num_cols)] for _ in range(num_rows)]

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell_xml = table_xml.xpath(
                f".//a:tbl/a:tr[{row_idx + 1}]/a:tc[{col_idx + 1}]"
            )
            if not cell_xml:
                continue

            cell_xml = cell_xml[0]
            row_span = int(cell_xml.get("rowSpan", 1))
            col_span = int(cell_xml.get("gridSpan", 1))

            # Place text in the grid
            # remove newline char to prevserve table structure
            cleaned_text = re.sub(r"[\n\r]", "", cell.text)
            grid[row_idx][col_idx] = cleaned_text

            # Mark spanned cells
            for i in range(row_span):
                for j in range(col_span):
                    if i == 0 and j == 0:
                        continue
                    if row_idx + i < num_rows and col_idx + j < num_cols:
                        grid[row_idx + i][col_idx + j] = ""

    # Convert grid to Markdown format
    table_text = []
    header = "|" + "|".join(grid[0]) + "|"
    separator = "|" + "---|" * num_cols
    table_text.append(header)
    table_text.append(separator)
    for row in grid[1:]:
        line = "|" + "|".join(row) + "|"
        table_text.append(line)
        # print(line)

    return table_text


def handle_grouped_shapes(shape) -> list[str]:
    """
    Formats grouped shapes into Markdown.
    """
    group_text = []

    def handle_shapes(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for grouped_shape in shape.shapes:
                handle_shapes(grouped_shape)
        else:
            if shape.has_text_frame:
                group_text.extend(handle_text_elements(shape))

    handle_shapes(shape)
    return group_text


def handle_charts(shape) -> list[str]:
    """
    Handles charts within a shape, converting them into Markdown format.
    """
    chart = shape.chart
    chart_title = chart.chart_title.text_frame.text if chart.has_title else "Chart"
    chart_text = [f" {chart_title}:"]
    for series in chart.series:
        series_text = f"Series '{series.name}'"
        chart_text.append(series_text)
    return chart_text


def handle_author_notes(slide) -> list[str]:
    notes = []
    if slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            notes.append("Speaker Notes:")
            notes.append(notes_text)
    return notes


# TODO :azure form reco to extract text from images
def handle_pictures(shape):
    pass

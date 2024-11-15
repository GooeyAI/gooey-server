import codecs
import csv
import datetime
import hashlib
import io
import mimetypes
import multiprocessing
import re
import tempfile
import typing
from functools import partial
from time import time

import gooey_gui as gui
import numpy as np
import requests
from django.db import transaction
from django.db.models import F
from django.utils import timezone
from furl import furl
from loguru import logger
from pydantic import BaseModel, Field

from app_users.models import AppUser
from daras_ai.image_input import (
    upload_file_from_bytes,
    safe_filename,
    get_mimetype_from_response,
)
from daras_ai_v2 import settings
from daras_ai_v2.asr import (
    AsrModels,
    run_asr,
    run_google_translate,
    download_youtube_to_wav,
)
from daras_ai_v2.azure_doc_extract import (
    table_arr_to_prompt_chunked,
    THEAD,
    azure_doc_extract_page_num,
)
from daras_ai_v2.doc_search_settings_widgets import (
    is_user_uploaded_url,
)
from daras_ai_v2.embedding_model import create_embeddings_cached, EmbeddingModels
from daras_ai_v2.exceptions import raise_for_status, call_cmd, UserError
from daras_ai_v2.functional import (
    flatmap_parallel,
    map_parallel,
    flatmap_parallel_ascompleted,
)
from daras_ai_v2.gdrive_downloader import (
    gdrive_download,
    is_gdrive_url,
    is_gdrive_presentation_url,
    url_to_gdrive_file_id,
    gdrive_metadata,
)
from daras_ai_v2.redis_cache import redis_lock
from daras_ai_v2.scraping_proxy import (
    get_scraping_proxy_cert_path,
    requests_scraping_kwargs,
    SCRAPING_PROXIES,
)
from daras_ai_v2.search_ref import (
    SearchReference,
    remove_quotes,
    generate_text_fragment_url,
)
from daras_ai_v2.text_splitter import text_splitter, Document
from embeddings.models import EmbeddedFile, EmbeddingsReference
from files.models import FileMetadata


class DocSearchRequest(BaseModel):
    search_query: str
    keyword_query: str | list[str] | None

    documents: list[str] | None

    max_references: int | None
    max_context_words: int | None
    scroll_jump: int | None

    doc_extract_url: str | None

    embedding_model: typing.Literal[tuple(e.name for e in EmbeddingModels)] | None
    dense_weight: float | None = Field(
        ge=0.0,
        le=1.0,
        title="Dense Embeddings Weightage",
        description="""
Weightage for dense vs sparse embeddings. `0` for sparse, `1` for dense, `0.5` for equal weight.
Generally speaking, dense embeddings excel at understanding the context of the query, whereas sparse vectors excel at keyword matches.
        """,
    )


def references_as_prompt(references: list[SearchReference], sep="\n\n") -> str:
    """
    Convert a list of references to a prompt containing the formatted search results.

    Args:
        references: list of references
        sep: separator between references in the prompt

    Returns:
        prompt string
    """
    return sep.join(
        f'''\
Search Result: [{idx + 1}]
Title: """{remove_quotes(ref["title"])}"""
Snippet: """
{remove_quotes(ref["snippet"])}
"""\
'''
        for idx, ref in enumerate(references)
    )


def get_top_k_references(
    request: DocSearchRequest, is_user_url: bool = True, current_user: AppUser = None
) -> typing.Generator[str, None, list[SearchReference]]:
    """
    Get the top k documents that ref the search query

    Args:
        request: the document search request
        is_user_url: whether the url is user-uploaded
        current_user: the current user

    Returns:
        the top k documents
    """
    from recipes.BulkRunner import url_to_runs

    yield "Fetching latest knowledge docs..."
    input_docs = request.documents or []

    if request.doc_extract_url:
        page_cls, sr, pr = url_to_runs(request.doc_extract_url)
        selected_asr_model = sr.state.get("selected_asr_model")
        google_translate_target = sr.state.get("google_translate_target")
    else:
        selected_asr_model = google_translate_target = None

    file_url_metas = flatmap_parallel(doc_or_yt_url_to_metadatas, input_docs)
    file_urls, file_metas = zip(*file_url_metas)

    yield "Creating knowledge embeddings..."

    embedding_model = EmbeddingModels.get(
        request.embedding_model,
        default=EmbeddingModels.get(
            EmbeddedFile._meta.get_field("embedding_model").default
        ),
    )
    embedded_files: list[EmbeddedFile] = map_parallel(
        lambda f_url, file_meta: get_or_create_embedded_file(
            f_url=f_url,
            file_meta=file_meta,
            max_context_words=request.max_context_words,
            scroll_jump=request.scroll_jump,
            google_translate_target=google_translate_target,
            selected_asr_model=selected_asr_model,
            embedding_model=embedding_model,
            is_user_url=is_user_url,
            current_user=current_user,
        ),
        file_urls,
        file_metas,
        max_workers=4,
    )
    if not embedded_files:
        yield "No embeddings found - skipping search"
        return []

    yield "Searching knowledge base..."

    vespa_file_ids = [ref.vespa_file_id for ref in embedded_files]
    EmbeddedFile.objects.filter(id__in=[ref.id for ref in embedded_files]).update(
        query_count=F("query_count") + 1,
        last_query_at=timezone.now(),
    )
    # chunk_count = sum(len(ref.document_ids) for ref in embedding_refs)
    # logger.debug(f"Knowledge base has {len(file_ids)} documents ({chunk_count} chunks)")

    s = time()
    search_result = query_vespa(
        request.search_query,
        file_ids=vespa_file_ids,
        limit=request.max_references or 100,
        embedding_model=embedding_model,
        semantic_weight=(
            request.dense_weight if request.dense_weight is not None else 1.0
        ),
    )
    references = list(vespa_search_results_to_refs(search_result))
    logger.debug(f"Search returned {len(references)} references in {time() - s:.2f}s")

    # merge duplicate references
    uniques: dict[str, SearchReference] = {}
    for ref in references:
        key = ref["url"]
        try:
            existing = uniques[key]
        except KeyError:
            uniques[key] = ref
        else:
            existing["snippet"] += "\n\n...\n\n" + ref["snippet"]
            existing["score"] = (existing["score"] + ref["score"]) / 2
    return list(uniques.values())


def vespa_search_results_to_refs(
    search_result: dict,
) -> typing.Iterable[SearchReference]:
    for hit in search_result["root"].get("children", []):
        try:
            ref = EmbeddingsReference.objects.get(vespa_doc_id=hit["fields"]["id"])
        except EmbeddingsReference.DoesNotExist:
            continue
        if "text/html" in ref.embedded_file.metadata.mime_type:
            # logger.debug(f"Generating fragments {ref['url']} as it is a HTML file")
            ref.url = generate_text_fragment_url(url=ref.url, text=ref.snippet)
        yield SearchReference(
            url=ref.url, title=ref.title, snippet=ref.snippet, score=hit["relevance"]
        )


def query_vespa(
    search_query: str,
    file_ids: list[str],
    limit: int,
    embedding_model: EmbeddingModels,
    semantic_weight: float = 1.0,
) -> dict:
    query_embedding = create_embeddings_cached([search_query], model=embedding_model)[0]
    if query_embedding is None or not file_ids:
        return {"root": {"children": []}}
    file_ids_str = ", ".join(map(repr, file_ids))
    query = f"select * from {settings.VESPA_SCHEMA} where file_id in (@fileIds) and (userQuery() or ({{targetHits: {limit}}}nearestNeighbor(embedding, q))) limit {limit}"
    logger.debug(f"Vespa query: {'-'*80}\n{query}\n{'-'*80}")
    if semantic_weight == 1.0:
        ranking = "semantic"
    elif semantic_weight == 0.0:
        ranking = "bm25"
    else:
        ranking = "fusion"
    response = get_vespa_app().query(
        yql=query,
        query=search_query,
        ranking=ranking,
        body={
            "ranking.features.query(q)": padded_embedding(query_embedding),
            "ranking.features.query(semanticWeight)": semantic_weight,
            "fileIds": file_ids_str,
        },
    )
    assert response.is_successful()
    return response.get_json()


def get_vespa_app():
    from vespa.application import Vespa

    return Vespa(url=settings.VESPA_URL)


def doc_or_yt_url_to_metadatas(f_url: str) -> list[tuple[str, FileMetadata]]:
    if is_yt_dlp_able_url(f_url):
        entries = yt_dlp_get_video_entries(f_url)
        return [
            (
                entry["webpage_url"],
                FileMetadata(
                    name=entry.get("title", "YouTube Video"),
                    # youtube doesn't provide etag, so we use filesize_approx or upload_date
                    etag=entry.get("filesize_approx") or entry.get("upload_date"),
                    # we will later convert & save as wav
                    mime_type="audio/wav",
                    total_bytes=entry.get("filesize_approx", 0),
                ),
            )
            for entry in entries
        ]
    else:
        return [(f_url, doc_url_to_file_metadata(f_url))]


def doc_url_to_file_metadata(f_url: str) -> FileMetadata:
    from googleapiclient.errors import HttpError

    f = furl(f_url.strip("/"))
    if is_gdrive_url(f):
        # extract filename from google drive metadata
        try:
            meta = gdrive_metadata(url_to_gdrive_file_id(f))
        except HttpError as e:
            if e.status_code == 404:
                raise UserError(
                    f"Could not download the google doc at {f_url} "
                    f"Please make sure to make the document public for viewing."
                ) from e
            else:
                raise
        name = meta["name"]
        etag = meta.get("md5Checksum") or meta.get("modifiedTime")
        mime_type = meta["mimeType"]
        total_bytes = int(meta.get("size") or 0)
    else:
        try:
            if is_user_uploaded_url(f_url):
                r = requests.head(f_url)
            else:
                r = requests.head(
                    f_url,
                    timeout=settings.EXTERNAL_REQUEST_TIMEOUT_SEC,
                    **requests_scraping_kwargs(),
                )
            raise_for_status(r)
        except requests.RequestException as e:
            logger.warning(f"ignore error while downloading {f_url}: {e}")
            name = None
            mime_type = None
            etag = None
            total_bytes = 0
        else:
            name = (
                r.headers.get("content-disposition", "")
                .split("filename=")[-1]
                .strip('"')
            )
            etag = r.headers.get("etag", r.headers.get("last-modified"))
            if etag:
                etag = etag.strip('"')
            mime_type = get_mimetype_from_response(r)
            total_bytes = int(r.headers.get("content-length") or 0)
    # extract filename from url as a fallback
    if not name:
        if is_user_uploaded_url(f_url):
            name = f.path.segments[-1]
        else:
            name = f"{f.host}{f.path}"
    # guess mimetype from name as a fallback
    if not mime_type:
        mime_type = mimetypes.guess_type(name)[0]
    return FileMetadata(
        name=name, etag=etag, mime_type=mime_type or "", total_bytes=total_bytes
    )


def yt_dlp_get_video_entries(url: str) -> list[dict]:
    data = yt_dlp_extract_info(url)
    entries = data.get("entries", [data])
    return [e for e in entries if e]


def yt_dlp_extract_info(url: str) -> dict:
    import yt_dlp

    # https://github.com/yt-dlp/yt-dlp/blob/master/yt_dlp/options.py
    params = dict(
        ignoreerrors=True,
        check_formats=False,
        proxy=SCRAPING_PROXIES.get("https"),
        client_certificate=get_scraping_proxy_cert_path(),
    )
    with yt_dlp.YoutubeDL(params) as ydl:
        data = ydl.extract_info(url, download=False)
        if not data:
            raise UserError(
                f"Could not download the youtube video at {url!r}. "
                f"Please make sure the video is public and the url is correct."
            )
        return data


def get_or_create_embedded_file(
    *,
    f_url: str,
    file_meta: FileMetadata,
    max_context_words: int,
    scroll_jump: int,
    google_translate_target: str | None,
    selected_asr_model: str | None,
    embedding_model: EmbeddingModels,
    is_user_url: bool,
    current_user: AppUser,
) -> EmbeddedFile:
    """
    Return Vespa document ids and document tags
    for a given document url + metadata.
    """
    lookup = dict(
        url=f_url,
        metadata__name=file_meta.name,
        metadata__etag=file_meta.etag,
        metadata__mime_type=file_meta.mime_type,
        metadata__total_bytes=file_meta.total_bytes,
        max_context_words=max_context_words,
        scroll_jump=scroll_jump,
        google_translate_target=google_translate_target or "",
        selected_asr_model=selected_asr_model or "",
        embedding_model=embedding_model.name,
    )
    file_id = hashlib.sha256(str(lookup).encode()).hexdigest()
    with redis_lock(f"gooey/get_or_create_embeddings/v1/{file_id}"):
        try:
            return EmbeddedFile.objects.filter(**lookup).order_by("-updated_at")[0]
        except IndexError:
            refs = create_embeddings_in_search_db(
                f_url=f_url,
                file_meta=file_meta,
                file_id=file_id,
                max_context_words=max_context_words,
                scroll_jump=scroll_jump,
                google_translate_target=google_translate_target or "",
                selected_asr_model=selected_asr_model or "",
                embedding_model=embedding_model,
                is_user_url=is_user_url,
            )
            with transaction.atomic():
                file_meta.save()
                embedded_file = EmbeddedFile.objects.get_or_create(
                    **lookup,
                    defaults=dict(
                        metadata=file_meta,
                        vespa_file_id=file_id,
                        created_by=current_user,
                    ),
                )[0]
                for ref in refs:
                    ref.embedded_file = embedded_file
                EmbeddingsReference.objects.bulk_create(refs)
            return embedded_file


def create_embeddings_in_search_db(
    *,
    f_url: str,
    file_meta: FileMetadata,
    file_id: str,
    max_context_words: int,
    scroll_jump: int,
    google_translate_target: str | None,
    selected_asr_model: str | None,
    embedding_model: EmbeddingModels,
    is_user_url: bool,
) -> list[EmbeddingsReference]:
    refs = []
    vespa = get_vespa_app()
    for ref, embedding in get_embeds_for_doc(
        f_url=f_url,
        file_meta=file_meta,
        max_context_words=max_context_words,
        scroll_jump=scroll_jump,
        google_translate_target=google_translate_target,
        selected_asr_model=selected_asr_model,
        embedding_model=embedding_model,
        is_user_url=is_user_url,
    ):
        doc_id = file_id + "/" + hashlib.sha256(str(ref).encode()).hexdigest()
        db_ref = EmbeddingsReference(
            vespa_doc_id=doc_id,
            url=ref["url"],
            title=ref["title"],
            snippet=ref["snippet"],
        )
        refs.append(db_ref)
        vespa.feed_data_point(
            schema=settings.VESPA_SCHEMA,
            data_id=doc_id,
            fields=format_embedding_row(
                doc_id=doc_id,
                created_at=db_ref.created_at,
                file_id=file_id,
                ref=ref,
                embedding=embedding,
            ),
            operation_type="feed",
        )
    return refs


def get_embeds_for_doc(
    *,
    f_url: str,
    file_meta: FileMetadata,
    max_context_words: int,
    scroll_jump: int,
    google_translate_target: str | None,
    selected_asr_model: str | None,
    embedding_model: EmbeddingModels,
    is_user_url: bool,
) -> typing.Iterator[tuple[SearchReference, np.ndarray]]:
    """
    Get document embeddings for a given document url.

    Args:
        f_url: document url
        file_meta: document metadata
        max_context_words: max number of words to include in each chunk
        scroll_jump: number of words to scroll by
        embedding_model: selected embedding model
        google_translate_target: target language for google translate
        selected_asr_model: selected ASR model (used for audio files)
        is_user_url: whether the url is user-uploaded

    Returns:
        list of (metadata, embeddings) tuples
    """
    pages = doc_url_to_text_pages(
        f_url=f_url,
        file_meta=file_meta,
        selected_asr_model=selected_asr_model,
        is_user_url=is_user_url,
    )
    refs = pages_to_split_refs(
        pages=pages,
        f_url=f_url,
        file_meta=file_meta,
        max_context_words=max_context_words,
        scroll_jump=scroll_jump,
    )
    translate_split_refs(refs, google_translate_target)
    texts = [m["title"] + " | " + m["snippet"] for m in refs]
    # get doc embeds in batches
    batch_size = 16  # azure openai limits
    return flatmap_parallel_ascompleted(
        partial(create_embeddings_cached, model=embedding_model),
        [refs[i : i + batch_size] for i in range(0, len(refs), batch_size)],
        [texts[i : i + batch_size] for i in range(0, len(texts), batch_size)],
        max_workers=2,
    )


def translate_split_refs(
    refs: list[SearchReference], google_translate_target: str | None
):
    if not google_translate_target:
        return
    snippets = [ref["snippet"] for ref in refs]
    translated_snippets = run_google_translate(snippets, google_translate_target)
    for ref, translated_snippet in zip(refs, translated_snippets):
        ref["snippet"] = translated_snippet


def pages_to_split_refs(
    *,
    pages,
    f_url: str,
    file_meta: FileMetadata,
    max_context_words: int,
    scroll_jump: int,
) -> list[SearchReference]:
    import pandas as pd

    chunk_size = int(max_context_words * 2)
    chunk_overlap = int(max_context_words * 2 / scroll_jump)
    if isinstance(pages, pd.DataFrame):
        refs = []
        # treat each row as a separate document
        for idx, row in pages.iterrows():
            row = dict(row)
            sections = (row.pop("sections", "") or "").strip()
            snippet = (row.pop("snippet", "") or "").strip()
            if sections:
                splits = split_sections(
                    sections, chunk_size=chunk_size, chunk_overlap=chunk_overlap
                )
            elif snippet:
                splits = text_splitter(
                    snippet, chunk_size=chunk_size, chunk_overlap=chunk_overlap
                )
            else:
                continue
            refs += [
                {
                    "title": file_meta.name,
                    "url": f_url,
                    **row,  # preserve extra csv rows
                    "snippet": doc.text,
                    **doc.kwargs,
                    "score": -1,
                }
                for doc in splits
            ]
    else:
        # split the text into chunks
        refs = [
            {
                "title": (
                    file_meta.name + (f", page {doc.end + 1}" if len(pages) > 1 else "")
                ),
                "url": add_page_number_to_pdf(
                    f_url, (doc.end + 1 if len(pages) > 1 else None)
                ).url,
                "snippet": doc.text,
                **doc.kwargs,
                "score": -1,
            }
            for doc in text_splitter(
                pages, chunk_size=chunk_size, chunk_overlap=chunk_overlap
            )
        ]
    return refs


sections_re = re.compile(r"(\s*[\r\n\f\v]|^)(\w+)\=", re.MULTILINE)


def split_sections(
    sections: str, *, chunk_overlap: int, chunk_size: int
) -> typing.Iterable[Document]:
    split = sections_re.split(sections)
    header = ""
    for i in range(2, len(split), 3):
        role = split[i].strip()
        content = split[i + 1].strip()
        if not content:
            continue
        match role:
            case "content" | "table":
                for doc in text_splitter(
                    content, chunk_size=chunk_size, chunk_overlap=chunk_overlap
                ):
                    doc.text = header + doc.text
                    yield doc

            case "csv":
                reader = csv.reader(io.StringIO(content))
                for prompt in table_arr_to_prompt_chunked(reader, chunk_size=1024):
                    yield Document(prompt, (0, 0))

            case _:
                header += f"{role}={content}\n"


def doc_url_to_text_pages(
    *,
    f_url: str,
    file_meta: FileMetadata,
    selected_asr_model: str | None,
    is_user_url: bool = True,
) -> typing.Union[list[str], "pd.DataFrame"]:
    """
    Download document from url and convert to text pages.
    """
    f_bytes, mime_type = download_content_bytes(
        f_url=f_url, mime_type=file_meta.mime_type, is_user_url=is_user_url
    )
    if not f_bytes:
        return []
    return any_bytes_to_text_pages_or_df(
        f_url=f_url,
        f_name=file_meta.name,
        f_bytes=f_bytes,
        mime_type=mime_type,
        selected_asr_model=selected_asr_model,
    )


def download_content_bytes(
    *, f_url: str, mime_type: str, is_user_url: bool = True
) -> tuple[bytes, str]:
    if is_yt_dlp_able_url(f_url):
        return download_youtube_to_wav(f_url), "audio/wav"
    f = furl(f_url)
    if is_gdrive_url(f):
        # download from google drive
        return gdrive_download(f, mime_type)
    try:
        # download from url
        if is_user_uploaded_url(f_url):
            r = requests.get(f_url)
        else:
            r = requests.get(f_url, **requests_scraping_kwargs())
        raise_for_status(r, is_user_url=is_user_url)
    except requests.RequestException as e:
        logger.warning(f"ignore error while downloading {f_url}: {e}")
        return b"", ""
    f_bytes = r.content
    # if it's a known encoding, standardize to utf-8
    encoding = r.headers.get("content-type", "").split("charset=")[-1].strip('"')
    if encoding:
        try:
            codec = codecs.lookup(encoding)
        except LookupError:
            pass
        else:
            try:
                f_bytes = codec.decode(f_bytes)[0].encode()
            except UnicodeDecodeError:
                pass
    mime_type = get_mimetype_from_response(r)
    return f_bytes, mime_type


def any_bytes_to_text_pages_or_df(
    *,
    f_url: str,
    f_name: str,
    f_bytes: bytes,
    mime_type: str,
    selected_asr_model: str | None,
) -> typing.Union[list[str], "pd.DataFrame"]:
    if mime_type.startswith("audio/") or mime_type.startswith("video/"):
        if is_gdrive_url(furl(f_url)) or is_yt_dlp_able_url(f_url):
            f_url = upload_file_from_bytes(f_name, f_bytes, content_type=mime_type)
        transcript = run_asr(
            f_url,
            selected_model=(selected_asr_model or AsrModels.whisper_large_v2.name),
        )
        return [transcript]

    try:
        return pdf_or_tabular_bytes_to_text_pages_or_df(
            f_url=f_url,
            f_name=f_name,
            f_bytes=f_bytes,
            mime_type=mime_type,
            # for now, only use form recognizer if asr model is selected.
            # We should later change the doc extract settings to include the form recognizer option
            use_form_reco=bool(selected_asr_model),
        )
    except UnsupportedDocumentError:
        pass

    if mime_type == "text/plain":
        text = f_bytes.decode()
    else:
        ext = mimetypes.guess_extension(mime_type) or ""
        text = pandoc_to_text(f_name + ext, f_bytes)
    return [text]


def is_yt_dlp_able_url(url: str) -> bool:
    f = furl(url)
    return (
        "youtube.com" in f.origin
        or "youtu.be" in f.origin
        or "fb.watch" in f.origin
        or (
            ("facebook.com" in f.origin or "fb.com" in f.origin)
            and (
                "videos" in f.path.segments
                or "/share/v/" in f.pathstr
                or "v" in f.query.params
            )
        )
    )


def pdf_or_tabular_bytes_to_text_pages_or_df(
    *,
    f_url: str,
    f_name: str,
    f_bytes: bytes,
    mime_type: str,
    use_form_reco: bool,
):
    import pandas as pd

    if mime_type == "application/pdf":
        if use_form_reco:
            df = pdf_to_form_reco_df(f_url, f_name, f_bytes, mime_type)
        else:
            return pdf_to_text_pages(f=io.BytesIO(f_bytes))

    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ]:
        if use_form_reco:
            return pptx_to_form_reco(f_url, f_name, f_bytes, mime_type)
        else:
            return pptx_to_text_pages(f=io.BytesIO(f_bytes))

    else:
        df = tabular_bytes_to_str_df(
            f_name=f_name, f_bytes=f_bytes, mime_type=mime_type
        )

    if "sections" in df.columns or "snippet" in df.columns:
        return df
    else:
        df.columns = [THEAD + col + THEAD for col in df.columns]
        return pd.DataFrame(["csv=" + df.to_csv(index=False)], columns=["sections"])


def tabular_bytes_to_str_df(
    *,
    f_name: str,
    f_bytes: bytes,
    mime_type: str,
) -> "pd.DataFrame":
    df = tabular_bytes_to_any_df(
        f_name=f_name, f_bytes=f_bytes, mime_type=mime_type, dtype=str
    )
    return df.fillna("")


def tabular_bytes_to_any_df(
    *,
    f_name: str,
    f_bytes: bytes,
    mime_type: str,
    dtype=None,
):
    import pandas as pd

    f = io.BytesIO(f_bytes)
    match mime_type:
        case "text/csv":
            df = pd.read_csv(f, dtype=dtype)
        case "text/tab-separated-values":
            df = pd.read_csv(f, sep="\t", dtype=dtype)
        case "application/json":
            df = pd.read_json(f, dtype=dtype)
        case "application/xml":
            df = pd.read_xml(f, dtype=dtype)
        case _ if "excel" in mime_type or "spreadsheet" in mime_type:
            df = pd.read_excel(f, dtype=dtype)
        case _:
            raise UnsupportedDocumentError(
                f"Unsupported document {mime_type=} ({f_name})"
            )
    return df

def pptx_to_text_pages(f: typing.BinaryIO) -> list[str]:
    """
    Extracts and converts text, tables, charts, and grouped shapes from a PPTX file into Markdown format.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(f)
    slides_text = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_content = [f"Slide {slide_idx}:"]  # Markdown heading for slide
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        slide_content.append(f"{text}")
                if shape.has_table:
                    slide_content.extend(pptx_format_table(shape.table))
                if shape.has_chart:
                    chart = shape.chart
                    chart_title = (
                        chart.chart_title.text_frame.text
                        if chart.has_title
                        else "Chart"
                    )
                    chart_text = [f"  {chart_title}:"]
                    for series in chart.series:
                        series_text = f"Series '{series.name}'"
                        chart_text.append(series_text)
                    slide_content.extend(chart_text)
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    slide_content.extend(pptx_format_grouped_shape(shape))
            except Exception as e:
                # Catch any exceptions and append to slide content which should be handled better
                slide_content.append(f"  Error processing shape: {e}")
                
        slides_text.append("\n".join(slide_content))
    return slides_text


def pptx_format_table(table) -> list[str]:
    """
    Formats a Shape-table into Markdown.
    """
    table_text = []
    if len(table.rows) == 0:
        return table_text

    header_row = [cell.text.strip() for cell in table.rows[0].cells]
    num_columns = len(header_row)

    table_text.append(pptx_gen_table_row(header_row))
    table_text.append(pptx_gen_table_row([":-:" for _ in header_row]))

    for row_idx in range(1, len(table.rows)):
        row = table.rows[row_idx]
        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
        if row_text:
            table_text.append(pptx_gen_table_row(row_text))

    return table_text


def pptx_gen_table_row(row: list[str]) -> str:
    """
    Generates a formatted table row for Markdown.
    """
    return "| " + " | ".join([c.replace("\n", "<br />") for c in row]) + " |"


def pptx_format_grouped_shape(group_shape) -> list[str]:
    """
    Recursively formats grouped shapes, extracting text from each.
    """
    group_text = []
    check_recursively_for_text(group_shape, group_text)
    return group_text

def check_recursively_for_text(shape, group_text):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            check_recursively_for_text(sub_shape, group_text)
    else:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                group_text.append(f"{text}")
                
                
class UnsupportedDocumentError(UserError):
    pass


def pdf_to_text_pages(f: typing.BinaryIO) -> list[str]:
    import pdftotext

    return list(pdftotext.PDF(f))


def pptx_to_form_reco(f_url: str, f_name: str, f_bytes: bytes, mime_type: str) -> list[str]:

    if is_gdrive_url(furl(f_url)):
        f_url = upload_file_from_bytes(f_name, f_bytes, content_type=mime_type)
    num_slides = get_pptx_num_slides(f_bytes)

    # logger.debug(f"pptx_to_form_reco: {f_url} {num_slides}")

    return list(map_parallel(
        lambda slide_num: 
        azure_doc_extract_page_num(f_url, slide_num, model_id="prebuilt-read"),
        range(1, num_slides + 1),
        max_workers=4,
    ))


def get_pptx_num_slides(f_bytes: bytes) -> int:
    from pptx import Presentation

    with tempfile.NamedTemporaryFile() as infile:
        infile.write(f_bytes)
        prs = Presentation(infile.name)
        return len(prs.slides)


def pdf_to_form_reco_df(
    f_url: str,
    f_name: str,
    f_bytes: bytes,
    mime_type: str,
) -> "pd.DataFrame":
    import pandas as pd

    if is_gdrive_url(furl(f_url)):
        f_url = upload_file_from_bytes(f_name, f_bytes, content_type=mime_type)
    num_pages = get_pdf_num_pages(f_bytes)
    return pd.DataFrame(
        map_parallel(
            lambda page_num: (
                add_page_number_to_pdf(f_url, page_num).url,
                f"{f_name}, page {page_num}",
                azure_doc_extract_page_num(f_url, page_num),
            ),
            range(1, num_pages + 1),
            max_workers=4,
        ),
        columns=["url", "title", "sections"],
    )


def get_pdf_num_pages(f_bytes: bytes) -> int:
    with tempfile.NamedTemporaryFile() as infile:
        infile.write(f_bytes)
        output = call_cmd("pdfinfo", infile.name).lower()
        for line in output.splitlines():
            if not line.startswith("pages:"):
                continue
            try:
                return int(line.split("pages:")[-1])
            except ValueError:
                raise ValueError(f"Unexpected PDF Info: {line}")


def add_page_number_to_pdf(url: str | furl, page_num: int) -> furl:

    # if it's a google drive presentation, add the slide number to the fragment
    if is_gdrive_presentation_url(furl(url)):
        return furl(url).set(fragment_args={"slide": page_num} if page_num else {})

    return furl(url).set(fragment_args={"page": page_num} if page_num else {})


# dont use more than 1GB of memory for pandoc in total
MAX_PANDOC_MEM_MB = 512
_pandoc_lock = multiprocessing.Semaphore(4)  # semaphore ensures max pandoc processes


def pandoc_to_text(f_name: str, f_bytes: bytes, to="plain") -> str:
    """
    Convert document to text using pandoc.

    Args:
        f_name: filename of document
        f_bytes: document bytes
        to: pandoc output format (default: plain)

    Returns:
        extracted text content of document
    """
    with (
        _pandoc_lock,
        tempfile.NamedTemporaryFile("wb", suffix="." + safe_filename(f_name)) as infile,
        tempfile.NamedTemporaryFile("r") as outfile,
    ):
        infile.write(f_bytes)
        call_cmd(
            "pandoc",
            # https://pandoc.org/MANUAL.html#a-note-on-security
            "+RTS", f"-M{MAX_PANDOC_MEM_MB}M", "-RTS", "--sandbox",
            "--standalone",
            infile.name,
            "--wrap", "none",
            "--to", to,
            "--output",
            outfile.name,
        )  # fmt: skip
        return outfile.read()


def render_sources_widget(refs: list[SearchReference]):
    if not refs:
        return
    with gui.expander("💁‍♀️ Sources"):
        for idx, ref in enumerate(refs):
            gui.html(
                # language=HTML
                f"""<p>{idx + 1}. <a href="{ref['url']}" target="_blank">{ref['title']}</a></p>""",
            )
            gui.text(ref["snippet"], style={"maxHeight": "200px"})
        gui.write(
            "---\n"
            + "```text\n"
            + "\n".join(f"[{idx + 1}] {ref['url']}" for idx, ref in enumerate(refs))
            + "\n```",
            disabled=True,
        )


def format_embedding_row(
    doc_id: str,
    file_id: str,
    ref: SearchReference,
    embedding: np.ndarray,
    created_at: datetime.datetime,
):
    return dict(
        id=doc_id,
        file_id=file_id,
        embedding=padded_embedding(embedding),
        created_at=int(created_at.timestamp() * 1000),
        # url=ref["url"].encode("unicode-escape").decode(),
        # title=ref["title"].encode("unicode-escape").decode(),
        # snippet=ref["snippet"].encode("unicode-escape").decode(),
    )


EMBEDDING_SIZE = 3072


def padded_embedding(
    arr: np.ndarray, max_len: int = EMBEDDING_SIZE, pad_value: float = 0.0
) -> list:
    return [pad_value] * (max_len - len(arr)) + arr.tolist()
